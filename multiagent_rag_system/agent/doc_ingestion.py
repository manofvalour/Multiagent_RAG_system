"""
Three-stage document ingestion:
  FileParser      — extract text (PDF/DOCX/PPTX/image OCR/text)
  HybridChunker   — semantic boundary detection + recursive size enforcement
  embed + index   — batch embed chunks, write to VectorStore
"""
from __future__ import annotations

import asyncio
import io
import re
import time
import uuid
from pathlib import Path
from typing import Optional

import numpy as np
from langchain_text_splitters import Language, RecursiveCharacterTextSplitter
from sklearn.metrics.pairwise import cosine_similarity

from ..src.logger.logger import GLOBAL_LOGGER as logger
from ..src.exception.custom_exception import MulitagentragException
from ..src.embedding.embedding import get_embedder
from ..src.database.vector_store import get_vector_store
from ..src.utils.config_loader import get_settings
from ..src.models.models import (
    ContentType,
    DocumentChunk,
    IngestRequest,
    IngestResponse,
)

settings = get_settings()

#Detecting the type of content
def detect_content_type(text: str) -> ContentType:
    lines = text.strip().splitlines()
    if not lines:
        return ContentType.PROSE
    md_sig   = sum(1 for l in lines if re.match(r"^#{1,6}\s", l) or l.startswith(("- ", "* ", "> ", "```")))
    code_sig = sum(1 for l in lines if re.match(r"^\s*(def |class |import |from |    )", l))
    total    = max(len(lines), 1)
    if code_sig / total > 0.15:
        return ContentType.CODE
    if md_sig / total > 0.10:
        return ContentType.MARKDOWN
    return ContentType.PROSE


#File parser
class FileParser:
    """
    Extracts raw text from binary file uploads.

    PDF  : pypdf (text-native) → EasyOCR / PyMuPDF (scanned fallback)
    DOCX : python-docx, preserves heading hierarchy
    PPTX : python-pptx, slides + speaker notes
    Image: EasyOCR → pytesseract fallback
    Text : direct UTF-8 decode
    """

    async def parse(self, content: bytes, filename: str) -> tuple[str, ContentType]:
        ext = Path(filename).suffix.lower()
        if ext == ".pdf":
            return await self._parse_pdf(content)
        elif ext == ".docx":
            return self._parse_docx(content), ContentType.DOCX
        elif ext == ".pptx":
            return self._parse_pptx(content), ContentType.PPTX
        elif ext in (".png", ".jpg", ".jpeg", ".webp", ".tiff"):
            return await self._parse_image(content), ContentType.IMAGE
        elif ext == ".md":
            return content.decode("utf-8", errors="replace"), ContentType.MARKDOWN
        else:
            text = content.decode("utf-8", errors="replace")
            return text, detect_content_type(text)

    async def _parse_pdf(self, content: bytes) -> tuple[str, ContentType]:
        pages_text: list[str] = []
        needs_ocr:  list[int] = []

        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(content))
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if len(text.strip()) >= 50:
                    pages_text.append(f"[Page {i+1}]\n{text.strip()}")
                else:
                    needs_ocr.append(i)
        except ImportError:
            needs_ocr = list(range(20))   # assume all pages need OCR

        if needs_ocr:
            try:
                import easyocr
                import fitz   # PyMuPDF
                ocr = easyocr.Reader(["en"], gpu=False)
                doc = fitz.open(stream=content, filetype="pdf")
                for page_num in needs_ocr:
                    pix  = doc[page_num].get_pixmap(dpi=200)
                    img  = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
                    results = ocr.readtext(img, detail=0)
                    ocr_text = " ".join(results).strip()
                    if ocr_text:
                        pages_text.append(f"[Page {page_num+1} OCR]\n{ocr_text}")
            except ImportError:
                logger.warning("easyocr/PyMuPDF not installed — scanned PDF pages skipped")

        return "\n\n".join(pages_text), ContentType.PDF

    def _parse_docx(self, content: bytes) -> str:
        try:
            from docx import Document
            doc   = Document(io.BytesIO(content))
            lines = []
            for para in doc.paragraphs:
                if not para.text.strip():
                    continue
                if para.style.name.startswith("Heading"):
                    level = para.style.name.split()[-1]
                    lines.append(f"{'#' * int(level)} {para.text}")
                else:
                    lines.append(para.text)
            return "\n\n".join(lines)
        except ImportError:
            return "[python-docx not installed — run: pip install python-docx]"

    def _parse_pptx(self, content: bytes) -> str:
        try:
            from pptx import Presentation
            prs    = Presentation(io.BytesIO(content))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        texts.append(f"[Notes: {notes}]")
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except ImportError:
            return "[python-pptx not installed — run: pip install python-pptx]"

    async def _parse_image(self, content: bytes) -> str:
        try:
            import easyocr
            from PIL import Image
            ocr = easyocr.Reader(["en"], gpu=False)
            img = np.array(Image.open(io.BytesIO(content)))
            return " ".join(ocr.readtext(img, detail=0))
        except ImportError:
            try:
                import pytesseract
                from PIL import Image
                return pytesseract.image_to_string(Image.open(io.BytesIO(content)))
            except ImportError:
                return "[OCR unavailable — install easyocr or pytesseract]"


#Semantic chunker
class SemanticChunker:
    """
    Splits text at points where cosine similarity between adjacent
    sentences drops below threshold — i.e. where the topic changes.
    Uses a buffer window to smooth noisy single-sentence embeddings.
    """

    def __init__(self, threshold: float = 0.5, buffer_size: int = 1) -> None:
        self.threshold   = threshold
        self.buffer_size = buffer_size

    def split(self, text: str, content_type: ContentType) -> list[str]:
        sentences = self._to_sentences(text, content_type)
        if len(sentences) <= 2:
            return [text]
        embeddings   = self._embed_buffered(sentences)
        split_points = self._find_splits(embeddings)
        return self._merge(sentences, split_points)

    def _to_sentences(self, text: str, ct: ContentType) -> list[str]:
        if ct == ContentType.PROSE:
            raw = re.split(r"(?<=[.!?])\s+", text)
        else:
            raw = re.split(r"\n{2,}", text)
        return [s.strip() for s in raw if s.strip()]

    def _embed_buffered(self, sentences: list[str]) -> np.ndarray:
        buffered = [
            " ".join(sentences[max(0, i - self.buffer_size): i + self.buffer_size + 1])
            for i in range(len(sentences))
        ]
        return get_embedder.embed(
            buffered, normalize_embeddings=True, show_progress_bar=False
        ).astype(np.float32)

    def _find_splits(self, embeddings: np.ndarray) -> list[int]:
        splits, seg_start = [], 0
        for i in range(len(embeddings) - 1):
            sim = float(cosine_similarity(
                embeddings[i].reshape(1, -1),
                embeddings[i + 1].reshape(1, -1),
            )[0][0])
            if sim < self.threshold and (i - seg_start + 1) >= 2:
                splits.append(i)
                seg_start = i + 1
        return splits

    def _merge(self, sentences: list[str], splits: list[int]) -> list[str]:
        split_set = set(splits)
        segs, cur = [], []
        for i, s in enumerate(sentences):
            cur.append(s)
            if i in split_set:
                segs.append(" ".join(cur))
                cur = []
        if cur:
            segs.append(" ".join(cur))
        return segs


#Hybrid chunker
class HybridChunker:
    """
    Stage 1 — SemanticChunker: split on topic boundaries
    Stage 2 — RecursiveCharacterTextSplitter: enforce max size
    Segments already within the size limit pass through stage 2 untouched,
    so semantic boundaries are preserved wherever possible.
    """

    def __init__(self) -> None:
        self.semantic = SemanticChunker(threshold=0.5)
        if settings.use_tokens:
            import tiktoken
            enc = tiktoken.encoding_for_model("openai/gpt-oss-120b")
            self._len_fn = lambda t: len(enc.encode(t))
        else:
            self._len_fn = len

    def _get_splitter(self, ct: ContentType) -> RecursiveCharacterTextSplitter:
        if ct == ContentType.CODE:
            return RecursiveCharacterTextSplitter.from_language(
                language=Language.PYTHON,
                chunk_size=settings.chunk_size,
                chunk_overlap=settings.chunk_overlap,
            )
        if ct in (ContentType.MARKDOWN, ContentType.DOCX, ContentType.PPTX):
            return RecursiveCharacterTextSplitter(
                separators=["\n## ", "\n### ", "\n#### ", "\n\n", "\n", " ", ""],
                chunk_size=settings.chunk_size,
                chunk_overlap=settings.chunk_overlap,
                length_function=self._len_fn,
            )
        return RecursiveCharacterTextSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
            length_function=self._len_fn,
            add_start_index=True,
        )

    def chunk(
        self, text: str, source: str,
        content_type: ContentType,
        metadata: dict,
    ) -> list[DocumentChunk]:
        splitter = self._get_splitter(content_type)
        segments = self.semantic.split(text, content_type)
        final: list[str] = []
        for seg in segments:
            if self._len_fn(seg) <= settings.chunk_size:
                final.append(seg)
            else:
                final.extend(splitter.split_text(seg))

        return [
            DocumentChunk(
                content=seg,
                source=source,
                chunk_index=i,
                content_type=content_type,
                metadata=metadata,
            )
            for i, seg in enumerate(final)
            if seg.strip()
        ]


#Ingestion pipeline
class DocumentIngestionPipeline:
    """
    Orchestrates: parse → chunk → embed → index.
    Accepts both IngestRequest (JSON/text) and raw binary files.
    """

    def __init__(self) -> None:
        self.parser       = FileParser()
        self.chunker      = HybridChunker()

    async def ingest_text(self, req: IngestRequest) -> IngestResponse:
        t0 = time.perf_counter()
        ct = detect_content_type(req.content)
        return await self._process(req.content, req.source, ct, req.metadata, t0)

    async def ingest_file(self, content: bytes, filename: str, metadata: dict = {}) -> IngestResponse:
        t0 = time.perf_counter()
        text, ct = await self.parser.parse(content, filename)
        return await self._process(text, filename, ct, metadata, t0)

    async def _process(
        self, text: str, source: str,
        ct: ContentType, metadata: dict, t0: float,
    ) -> IngestResponse:
        doc_id = str(uuid.uuid4())
        chunks = self.chunker.chunk(text, source, ct, {**metadata, "doc_id": doc_id})

        if not chunks:
            logger.info(f"No chunks produced for {source!r}")
            return IngestResponse(
                document_id=doc_id, source=source, num_chunks=0,
                content_type=ct.value, latency_ms=0.0,
            )

        loop = asyncio.get_event_loop()
        embeddings = await loop.run_in_executor(
            None,
            lambda: get_embedder.embed(
                [c.content for c in chunks],
                normalize_embeddings=True,
                batch_size=64,
            ).astype(np.float32),
        )

        for chunk, emb in zip(chunks, embeddings):
            chunk.embedding = emb.tolist()

        await get_vector_store.add_chunks(chunks, embeddings)
        latency = round((time.perf_counter() - t0) * 1000, 2)
        logger.info(f"Ingested {source!r}  chunks={len(chunks)}  type={ct.value}  {latency}ms")

        return IngestResponse(
            document_id=doc_id,
            source=source,
            chunks_created=len(chunks),
            content_type=ct.value,
            latency_ms=latency,
        )